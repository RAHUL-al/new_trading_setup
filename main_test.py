from datetime import timezone
from google.cloud.firestore_v1.base_query import FieldFilter
from schemas import UserCreate, PasswordResetRequest, ForgotPasswordRequest, UserOut, UserUpdate, Token, LoginRequest, ForgotPasswordRequest, GeminiRequest, ChatRequest
from fastapi import Body
from fastapi.websockets import WebSocketState
import asyncio
from langchain_core.prompts import PromptTemplate
from langchain_ollama import ChatOllama
import aiohttp
from bs4 import BeautifulSoup
from langchain_core.messages import AIMessage, SystemMessage, HumanMessage
from langchain_core.output_parsers import StrOutputParser
from collections import Counter
import ujson as json
import redis
from datetime import datetime, timedelta
from fastapi import FastAPI, Depends, HTTPException, status, Request, WebSocket, WebSocketDisconnect, UploadFile, File, Form, Query
from fastapi.middleware.cors import CORSMiddleware
from pydub import AudioSegment
from dotenv import load_dotenv
from urllib.parse import parse_qs, unquote
from concurrent.futures import ThreadPoolExecutor
from fastapi.responses import FileResponse, JSONResponse
import time
from pdf2image import convert_from_path
from typing import List
from PIL import Image
import google.generativeai as genai
from pinecone import Pinecone
from auth import hash_password, verify_password, create_access_token
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Pinecone as PineconeVectorStore
import torch
import logging
from urllib.parse import parse_qs,unquote
from jose import JWTError, jwt
from ai_speech_module import Topic, AdvancedAudioProcessor
from langchain_ollama import OllamaLLM
import re
from typing import List
from fastapi import FastAPI, UploadFile, File, Form
from pdf2image import convert_from_path
from PIL import Image
from docx2pdf import convert as docx_to_pdf
import os, uuid, shutil, logging, tempfile
from langchain.text_splitter import RecursiveCharacterTextSplitter
from fastapi import HTTPException, status
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain_community.vectorstores import Pinecone as PineconeVectorStore
from PIL import Image, ImageDraw, ImageFont
import pandas as pd
import docx2txt
from pptx import Presentation
from firestore_models import FirestoreEssay
from firebase import db
from concurrent.futures import ProcessPoolExecutor
process_pool = ProcessPoolExecutor()
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, BackgroundTasks, WebSocket, WebSocketDisconnect
import aiohttp
import tempfile
from urllib.parse import parse_qs
from typing import List, Dict
from PIL import Image, ImageDraw, ImageFont
from pdf2image import convert_from_path
import docx2txt
from pptx import Presentation
import pandas as pd
from concurrent.futures import ThreadPoolExecutor, as_completed
import subprocess
from pdf2image import convert_from_path
from fastapi import HTTPException
from pathlib import Path
from PIL import Image
from typing import List, TypedDict, Annotated
import glob
from langchain.chains import LLMChain
import numpy as np
from scipy import signal
import librosa
from urllib.parse import parse_qs

python311_path = "/home/ubuntu/ai_speech/venv/bin/python"
script_path = "/home/ubuntu/ai_speech/ai_speech_model/ai_speech_fastapi/layoutparser_file.py"


CPU_API_BASE = "http://13.200.201.10:8000"
model_name = OllamaLLM(model="mistral")
scraping_api_key = os.getenv("SCRAPINGDOG_API_KEY")

PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_ENV = os.getenv("PINECONE_ENVIRONMENT")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
EMBEDDING_MODEL_NAME = "embaas/sentence-transformers-e5-large-v2"

SECRET_KEY = "jwt_secret_key"
ALGORITHM="HS256"


logging.basicConfig(
    filename='rag_log.log',
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

load_dotenv()
app = FastAPI(title="FastAPI‑Firebase")

redis_client = redis.StrictRedis(
    host=os.getenv("REDIS_HOST"),
    port=int(os.getenv("REDIS_PORT")),
    username=os.getenv("REDIS_USERNAME"),
    password=os.getenv("REDIS_PASSWORD"),
    decode_responses=True
)

origins = ["https://llm.edusmartai.com","http://localhost:3000","http://localhost:5173"]

app.add_middleware(CORSMiddleware, allow_origins=origins, allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

def get_user_from_redis_session(request: Request):
    token = request.headers.get("Authorization")
    if not token or not token.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid token")
    token = token.split(" ")[1]
    session_data = redis_client.get(f"session:{token}")
    if not session_data:
        raise HTTPException(status_code=401, detail="Session expired or invalid")
    return json.loads(session_data)

@app.post("/register", response_model=UserOut)
def register(user: UserCreate):
    user_ref = db.collection("users").where("username", "==", user.username).stream()
    if any(user_ref):
        raise HTTPException(400, "Username or email already exists")

    doc_ref = db.collection("users").add({
        "username": user.username,
        "email": user.email,
        "password": hash_password(user.password)
    })
    user_id = doc_ref[1].id
    return UserOut(id=user_id, username=user.username, email=user.email)

from fastapi.concurrency import run_in_threadpool
from fastapi import BackgroundTasks
from fastapi_mail import FastMail, MessageSchema, ConnectionConfig

PASSWORD_RESET_TOKEN_EXPIRE_HOURS = 1

mail_conf = ConnectionConfig(
    MAIL_USERNAME=os.getenv("MAIL_USERNAME", "testampli2023@gmail.com"),
    MAIL_PASSWORD=os.getenv("MAIL_PASSWORD", "mulpeeeuolzidejx"),
    MAIL_FROM=os.getenv("MAIL_FROM", "testampli2023@gmail.com"),
    MAIL_PORT=int(os.getenv("MAIL_PORT", 587)),
    MAIL_SERVER=os.getenv("MAIL_SERVER", "smtp.gmail.com"),
    MAIL_STARTTLS=bool(os.getenv("MAIL_STARTTLS", True)),
    MAIL_SSL_TLS=bool(os.getenv("MAIL_SSL_TLS", False)),
    USE_CREDENTIALS=True,
    VALIDATE_CERTS=True
)


async def validate_reset_token(token: str):
    try:
        token_data = redis_client.get(f"session:{token}")

        logging.info(f"token data {token_data}")

        if not token_data:
            raise HTTPException(status_code=400, detail="Invalid or expired token")
            
        return {"valid": True, "username": json.loads(token_data)["username"]}
        
    except HTTPException as he:
        raise he
    except Exception as e:
        logging.error(f"Error validating token: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Invalid or expired token"
        )


def generate_reset_token(user_id: str):
    """Generate JWT token for password reset"""
    expires = datetime.now(timezone.utc) + timedelta(hours=PASSWORD_RESET_TOKEN_EXPIRE_HOURS)
    payload = {
        "sub": user_id,
        "exp": expires,
        "type": "password_reset"
    }
    return jwt.encode(payload, SECRET_KEY, algorithm=ALGORITHM)



async def send_reset_email(email: str, token: str, background_tasks: BackgroundTasks):
    reset_link = f"https://llm.edusmartai.com/reset-password?token={token}"
    message = MessageSchema(
        subject="Reset Your Password",
        recipients=[email],
        body=f"Click the link to reset your password: {reset_link}",
        subtype="plain"
    )
    fm = FastMail(mail_conf)
    background_tasks.add_task(fm.send_message, message)


@app.post("/forgot-password")
async def forgot_password(request: ForgotPasswordRequest, background_tasks: BackgroundTasks):
    try:
        if not request.email or "@" not in request.email:
            return {"detail": "If this email exists in our system, you'll receive a password reset link"}

        docs = db.collection("users").where(filter=FieldFilter("email", "==", request.email)).stream()
        user_doc = next(docs, None)

        if not user_doc:
            logging.info(f"Password reset requested for non-existent email: {request.email}")
            return {"detail": "If this email exists in our system, you'll receive a password reset link"}

        reset_token = generate_reset_token(user_doc.id)
        token_expiry = datetime.now(timezone.utc) + timedelta(hours=PASSWORD_RESET_TOKEN_EXPIRE_HOURS)

        redis_client.setex(
            f"reset_token:{reset_token}",
            int(timedelta(hours=PASSWORD_RESET_TOKEN_EXPIRE_HOURS).total_seconds()),
            json.dumps({
                "user_id": user_doc.id,
                "email": request.email,
                "expires_at": token_expiry.isoformat()
            })
        )

        await send_reset_email(request.email, reset_token, background_tasks)
        logging.info(f"Password reset token generated for {request.email}")

        return {"detail": "If this email exists in our system, you'll receive a password reset link"}

    except Exception as e:
        logging.error(f"Error in forgot-password for {request.email}: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="An error occurred while processing your request"
        )


@app.get("/validate-reset-token/{token}")
async def validate_reset_token_endpoint(token: str):
    try:
        token_data = redis_client.get(f"session:{token}")

        logging.info(f"token data {token_data}")

        if not token_data:
            raise HTTPException(status_code=400, detail="Invalid or expired token")
            
        return {"valid": True, "username": json.loads(token_data)["username"]}
        
    except HTTPException as he:
        raise he
    except Exception as e:
        logging.error(f"Error validating token: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Invalid or expired token"
        )


@app.post("/reset-password")
async def reset_password(request: PasswordResetRequest = Body(...)):
    try:
        token_data_raw = redis_client.get(f"session:{request.token}")
        if not token_data_raw:
            raise HTTPException(status_code=400, detail="Invalid or expired token")

        logging.info(f"token data {token_data_raw}")

        token_data = json.loads(token_data_raw)
        user_id = token_data["user_id"]

        user_ref = db.collection("users").document(user_id)
        user_doc = user_ref.get()
        if not user_doc.exists:
            raise HTTPException(404, "User not found")

        user_ref.update({
            "password": hash_password(request.new_password),
            "updated_at": datetime.now(timezone.utc)
        })

        redis_client.delete(f"session:{request.token}")

        logging.info(f"Password reset successful for user {user_id}")
        return {"detail": "Password has been reset successfully"}

    except HTTPException as he:
        raise he
    except Exception as e:
        logging.error(f"Reset password error: {str(e)}", exc_info=True)
        raise HTTPException(500, "Failed to reset password")



@app.post("/login", response_model=Token)
def login(data: LoginRequest):
    docs = db.collection("users").where("username", "==", data.username).stream()
    user_doc = next(docs, None)
    if not user_doc or not verify_password(data.password, user_doc.to_dict()["password"]):
        raise HTTPException(status.HTTP_401_UNAUTHORIZED, detail="Bad credentials")

    token = create_access_token({"sub": user_doc.id})
    logging.info(f"token in login file : {token}")
    redis_client.setex(f"session:{token}", timedelta(hours=1), json.dumps({"user_id": user_doc.id, "username": data.username}))
    return Token(access_token=token, username=data.username)



@app.get("/logout")
def logout(request: Request):
    token = request.headers.get("Authorization")
    if token and token.startswith("Bearer "):
        redis_client.delete(f"session:{token.split(' ')[1]}")
    return {"detail": "Logged out"}

@app.get("/me", response_model=UserOut)
def me(user=Depends(get_user_from_redis_session)):
    doc = db.collection("users").document(user["user_id"]).get()
    if not doc.exists:
        raise HTTPException(404, "User not found")
    data = doc.to_dict()
    return UserOut(id=doc.id, username=data["username"], email=data["email"])



@app.post("/generate-prompt")
async def generate_prompt(data: GeminiRequest, user=Depends(get_user_from_redis_session)):
    
    url = f"https://en.wikipedia.org/wiki/{data.topic}"
    api_endpoint = f"https://api.scrapingdog.com/scrape?api_key={scraping_api_key}&url={url}"

    text = ""
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(api_endpoint) as response:
                if response.status == 200:
                    html = await response.text()
                    soup = BeautifulSoup(html, "html.parser")
                    for script in soup(["script", "style", "noscript"]):
                        script.decompose()
                    text = soup.get_text(separator="\n", strip=True)
                else:
                    logging.info(f"Error: {response.status} - {await response.text()}")
    except Exception as e:
        logging.exception(f"Failed to scrape data: {e}")

    prompt = (
        f"Generate a essay for a student in class {data.student_class} with a {data.accent} accent, "
        f"on the topic '{data.topic}', and the mood is '{data.mood}' and give me essay should be less than 400 words "
        f"and in response did not want \n\n or \n and also not want word count thanks you this type of stuff and used {text} "
        f"content for as updated data from internet and which is helpful in created essay and please give me output in paragraph format only not in points."
    )

    username = user.get("username")
    topic = Topic()
    response_text = await topic.topic_data_model_for_Qwen(username, prompt)

    essay_data = FirestoreEssay(
        username=username,
        user_id=user["user_id"],
        student_class=data.student_class,
        accent=data.accent,
        topic=data.topic,
        mood=data.mood,
        content=response_text
    )

    write_time, doc_ref = db.collection("essays").add(essay_data.to_dict())
    essay_id = doc_ref.id

    return JSONResponse(content={
        "response": response_text,
        "essay_id": essay_id
    })




@app.get("/overall-scoring-by-id")
async def overall_scoring_by_listening_module(essay_id: str):
    topic = Topic()
    result = await topic.overall_scoring_by_listening_module(essay_id)
    return result

TEMP_DIR = os.path.abspath("audio_folder")
os.makedirs(TEMP_DIR, exist_ok=True)


@app.get("/overall-scoring-by-id-speech-module")
async def overall_scoring_by_speech_module(essay_id: str):
    topic = Topic()
    result = await topic.overall_scoring_by_speech_module(essay_id)
    return result



TEMP_DIR = os.path.abspath("audio_folder")
os.makedirs(TEMP_DIR, exist_ok=True)

@app.websocket("/ws/audio")
async def audio_ws(websocket: WebSocket):
    await websocket.accept()
    query_params = parse_qs(websocket.url.query)
    username = query_params.get("username", [None])[0]
    token = query_params.get("token", [None])[0]

    if not username or not token:
        await websocket.close(code=4001)
        logging.info("Username or token missing.")
        return
        

    logging.info(f"[WS] Authenticated connection from {username}")
    chunk_index = 0
    chunk_results = []
    text_output = []

    date_str = datetime.now().strftime("%Y-%m-%d")
    user_dir = os.path.join(TEMP_DIR, username, date_str)
    os.makedirs(user_dir, exist_ok=True)

    final_output = os.path.join(user_dir, f"{username}_output.wav")
    transcript_path = os.path.join(user_dir, f"{username}_transcript.txt")

    if os.path.exists(final_output):
        os.remove(final_output)
    if os.path.exists(transcript_path):
        os.remove(transcript_path)

    loop = asyncio.get_event_loop()

    try:
        topic = Topic()
        while True:
            message = await websocket.receive()

            if message["type"] == "websocket.disconnect":
                print(f"[WS] {username} disconnected.")
                break

            if message["type"] == "websocket.receive" and "bytes" in message:
                chunk_filename = os.path.join(user_dir, f"chunk_{chunk_index}.wav")
                audio = AudioSegment(
                    data=message["bytes"],
                    sample_width=2,
                    frame_rate=16000,
                    channels=1
                )
                audio.export(chunk_filename, format="wav")

                async with aiohttp.ClientSession() as session:
                    with open(chunk_filename, "rb") as f:
                        form = aiohttp.FormData()
                        form.add_field("file", f, filename=os.path.basename(chunk_filename), content_type="audio/wav")
                        async with session.post(f"{CPU_API_BASE}/detect-emotion", data=form) as res:
                            emotion_data = await res.json()
                            emotion = emotion_data.get("emotion")

                    transcribed_text = await topic.speech_to_text(chunk_filename, username)

                    # grammar = await topic.grammar_checking(transcribed_text)

                    async with session.post(
                        f"{CPU_API_BASE}/fluency-score",
                        json={"text": transcribed_text}
                    ) as res:
                        fluency_data = await res.json()
                        fluency = fluency_data.get("fluency")

                    with open(chunk_filename, "rb") as f:
                        form = aiohttp.FormData()
                        form.add_field("file", f, filename=os.path.basename(chunk_filename), content_type="audio/wav")
                        async with session.post(f"{CPU_API_BASE}/pronunciation-score", data=form) as res:
                            pron_data = await res.json()
                            pronunciation = pron_data.get("pronunciation")

                topic.update_realtime_stats(fluency, pronunciation, emotion)

                text_output.append(transcribed_text)

                chunk_result = {
                    "chunk_index": chunk_index,
                    "text": transcribed_text,
                    "emotion": emotion,
                    "fluency": fluency,
                    "pronunciation": pronunciation,
                    # "grammar":grammar,
                    # "silvero": silvero,
                    "file_path": chunk_filename
                }

                logging.info(f"[Chunk {chunk_index}] {chunk_result}")
                chunk_results.append(chunk_result)
                chunk_index += 1

    except WebSocketDisconnect:
        logging.warning(f"[WS] {username} forcibly disconnected.")

    finally:
        await loop.run_in_executor(None, merge_chunks, chunk_results, final_output)

        with open(transcript_path, "w", encoding="utf-8") as f:
            f.write(" ".join(text_output).strip())

        try:
            essays_ref = db.collection("essays").where("username", "==", username)
            essays = essays_ref.stream()
            today = datetime.now().date()
            latest_essay = max((doc for doc in essays if doc.create_time.date() == today), key=lambda d: d.create_time, default=None)

            if latest_essay:
                essay_ref = db.collection("essays").document(latest_essay.id)
                average_scores = topic.get_average_realtime_scores()
                essay_ref.update({
                    "chunks": chunk_results,
                    "average_scores": average_scores
                })
                logging.info(f"Updated essay {latest_essay.id}")

        except Exception as e:
            logging.error(f"[Firestore Update Error] {e}")

        for file in os.listdir(user_dir):
            if file.startswith("chunk_") and file.endswith(".wav"):
                try:
                    os.remove(os.path.join(user_dir, file))
                except Exception as e:
                    logging.warning(f"Failed to remove {file}: {e}")


def merge_chunks(chunk_files, final_output):
    logging.info("[Merge] Merging audio chunks...")
    combined = AudioSegment.empty()

    for chunk in chunk_files:
        file_path = chunk.get("file_path")
        if file_path and os.path.exists(file_path):
            audio = AudioSegment.from_file(file_path, format="wav")
            combined += audio
        else:
            logging.warning(f"[Merge] Skipping missing or invalid file: {file_path}")

    combined.export(final_output, format="wav")
    logging.info("[Merge] Merged audio file saved.")




@app.get("/get-tts-audio")
def get_tts_audio(username: str):
    folder = os.path.join("text_to_speech_audio_folder", username)
    file_path = os.path.join(folder, f"{username}_output.wav")

    timeout = 60
    poll_interval = 2
    waited = 0

    while waited < timeout:
        if os.path.exists(file_path):
            return FileResponse(file_path, media_type="audio/wav", filename=f"{username}_output.wav")
        time.sleep(poll_interval)
        waited += poll_interval

    raise HTTPException(status_code=408, detail="Audio file not generated within 1 minute.")



SUPPORTED_IMAGE_FORMATS = [".png", ".jpg", ".jpeg"]
SUPPORTED_TEXT_FORMATS = [".txt"]
SUPPORTED_PDF_FORMATS = [".pdf"]
SUPPORTED_DOC_FORMATS = [".docx"]
SUPPORTED_PPT_FORMATS = [".pptx"]
SUPPORTED_XLS_FORMATS = [".xlsx"]


genai.configure(api_key=GOOGLE_API_KEY)
gemini_model = genai.GenerativeModel("gemini-1.5-pro")
pc = Pinecone(api_key=PINECONE_API_KEY)

if PINECONE_INDEX_NAME not in pc.list_indexes().names():
    pc.create_index(
        name=PINECONE_INDEX_NAME,
        dimension=1024,
        metric="cosine",
        pods=1,
        pod_type="p1.x1"
    )
    logging.info(f"Created new Pinecone index {PINECONE_INDEX_NAME} with dimension 1024")

index = pc.Index(PINECONE_INDEX_NAME)

embedding_model = HuggingFaceEmbeddings(
    model_name=EMBEDDING_MODEL_NAME,
    model_kwargs={'device': 'cuda' if torch.cuda.is_available() else 'cpu'}
)

executor = ThreadPoolExecutor(max_workers=8)

async def run_in_threadpool(func, *args, **kwargs):
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(executor, lambda: func(*args, **kwargs))


def render_text_to_image(text: str, width=800, font_size=18) -> Image.Image:
    font = ImageFont.load_default()
    lines = []
    dummy_img = Image.new("RGB", (width, 1000))
    draw = ImageDraw.Draw(dummy_img)

    words = text.split()
    line = ""
    for word in words:
        test_line = f"{line} {word}".strip()
        bbox = draw.textbbox((0, 0), test_line, font=font)
        w = bbox[2] - bbox[0]
        if w < width - 40:
            line = test_line
        else:
            lines.append(line)
            line = word
    lines.append(line)

    height = font_size * len(lines) + 50
    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)
    y = 20
    for line in lines:
        draw.text((20, y), line, font=font, fill="black")
        y += font_size
    return img

def file_to_images(file_path: str) -> List[Image.Image]:
    ext = os.path.splitext(file_path)[1].lower()
    images = []

    if ext in SUPPORTED_PDF_FORMATS:
        with ThreadPoolExecutor() as pdf_executor:
            futures = []
            chunk_size = 10
            images = convert_from_path(file_path, dpi=200, thread_count=4)

    elif ext in SUPPORTED_IMAGE_FORMATS:
        images = [Image.open(file_path)]

    elif ext in SUPPORTED_DOC_FORMATS:
        try:
            text = docx2txt.process(file_path)
            img = render_text_to_image(text)
            images = [img]
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"DOCX to image failed: {e}")

    elif ext in SUPPORTED_PPT_FORMATS:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            img = render_text_to_image(text)
            images.append(img)

    elif ext in SUPPORTED_XLS_FORMATS:
        try:
            excel = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, df in excel.items():
                text = f"Sheet: {sheet_name}\n{df.to_string(index=False)}"
                img = render_text_to_image(text)
                images.append(img)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"XLSX to image failed: {e}")

    elif ext in SUPPORTED_TEXT_FORMATS:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        img = render_text_to_image(content)
        images = [img]

    else:
        raise HTTPException(status_code=400, detail="Unsupported file format")

    return images


async def process_single_image(filename,username,image: Image.Image, idx: int,image_path) -> str:
    try:
        today_date = datetime.now().strftime("%Y-%m-%d")
        base_dir = f"diagrams/{username}_{today_date}?extracted_images"
        
        response = await run_in_threadpool(
            gemini_model.generate_content,
            [
                image,
                "1. Extract all visible text from this image.\n"
                "2. If the image contains diagrams, illustrations, or charts, summarize what they represent in 2-3 lines.\n"
                "Return the result with clear separation: first the text, then the image summary (if any)."
            ]
        )
        text = response.text.strip()

        result = subprocess.run(
            [python311_path, script_path, username, image_path, str(idx), filename],
            capture_output=True,
            text=True
        )
        try:
            output = json.loads(result.stdout)
            print(f"Page {idx + 1} Output:", output["result"])
        except json.JSONDecodeError:
            print(f"Error decoding output for page {idx + 1}:", result.stdout)

        if result.stderr:
            print(f"Error from script on page {idx + 1}:", result.stderr)

        return f"\n\n--- Page/Image {idx + 1} ---\n{text}"
    except Exception as e:
        logging.error(f"OCR failed on image {idx + 1}: {e}")
        return f"\n\n--- Page/Image {idx + 1} FAILED ---"

async def extract_text_parallel(filename,username,file_path: str, timeout_per_page: int = 240) -> str:
    images = await run_in_threadpool(file_to_images, file_path)
    if not images:
        return ""
    
    all_text = ""
    tasks = []

    filename = filename.replace(" ","_")
    filename = filename.split(".")[0]
    output_image_dir = f"/home/ubuntu/ai_speech/ai_speech_model/ai_speech_fastapi/temp_images/{username}/{filename}"
    if os.path.exists(output_image_dir):
        if os.listdir(output_image_dir):
            shutil.rmtree(output_image_dir)

    os.makedirs(output_image_dir, exist_ok=True)

    
    
    for idx, image in enumerate(images):
        image_path = os.path.join(output_image_dir, f"page_{idx + 1}.jpg")
        image.save(image_path, "JPEG")
        task = asyncio.create_task(
            asyncio.wait_for(
                process_single_image(filename,username,image, idx,image_path),
                timeout=timeout_per_page
            )
        )
        tasks.append(task)
    
    for i, task in enumerate(asyncio.as_completed(tasks)):
        try:
            page_text = await task
            all_text += page_text
        except asyncio.TimeoutError:
            logging.warning(f"Timeout processing page {i + 1}")
            all_text += f"\n\n--- Page/Image {i + 1} TIMEOUT ---"
        except Exception as e:
            logging.error(f"Error processing page {i + 1}: {e}")
            all_text += f"\n\n--- Page/Image {i + 1} ERROR ---"
    
    return all_text

async def extract_text_with_retry(filename,username,file_path: str, timeout=240, max_retries=3) -> tuple[str, bool]:
    last_error = None
    timeout_occurred = False
    
    for attempt in range(1, max_retries + 2):
        try:
            task_start = time.time()
            
            task = asyncio.create_task(extract_text_parallel(filename,username,file_path))
            
            try:
                result = await asyncio.wait_for(task, timeout=timeout)
                elapsed = time.time() - task_start
                logging.info(f"OCR succeeded in attempt {attempt} ({elapsed:.2f}s)")
                return result, timeout_occurred
                
            except asyncio.TimeoutError:
                task.cancel()
                elapsed = time.time() - task_start
                logging.warning(
                    f"OCR timeout in attempt {attempt} after {elapsed:.2f}s "
                    f"(Timeout setting: {timeout}s)"
                )
                last_error = f"Timeout after {timeout}s"
                timeout_occurred = True
                
                if attempt <= max_retries:
                    backoff = min(2 ** attempt, 10)
                    await asyncio.sleep(backoff)
                    
        except Exception as e:
            last_error = str(e)
            logging.error(f"OCR attempt {attempt} failed: {last_error}")
            if attempt <= max_retries:
                await asyncio.sleep(1)

    raise HTTPException(
        status_code=504,
        detail=f"OCR failed after {max_retries + 1} attempts. Last error: {last_error}"
    )

@app.post("/upload/")
async def upload_file(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    student_class: str = Form(...),
    subject: str = Form(...),
    curriculum: str = Form(...),
    username : str = Form(...),
    
):

    file_path = None
    try:
        start_time = time.time()
        folder = f"uploads/{curriculum}/{student_class}/{subject}"
        os.makedirs(folder, exist_ok=True)
        file_path = os.path.join(folder, file.filename)

        filename = file.filename

        with open(file_path, "wb") as buffer:
            while chunk := await file.read(1024 * 1024):
                buffer.write(chunk)
        logging.info(f"File saved in {time.time()-start_time:.2f}s")

        extract_start = time.time()
        try:
            extracted_text, timeout_occurred = await extract_text_with_retry(
                filename,
                username,
                file_path,
                timeout=240,
                max_retries=2,
            )
            logging.info(f"Text extracted in {time.time()-extract_start:.2f}s")
            
            if not extracted_text.strip():
                raise HTTPException(
                    status_code=422,
                    detail="No text could be extracted from the file"
                )
                
        except HTTPException as e:
            raise e
        except Exception as e:
            logging.error(f"OCR processing failed: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail="Failed to process document content"
            )

        namespace = f"{curriculum}_{student_class}_{subject}"
        
        existing_entries = index.query(
            vector=embedding_model.embed_query(extracted_text[:1000]),
            top_k=1,
            filter={
                "filename": {"$eq": file.filename},
                "type": {"$eq": "document"}
            },
            namespace=namespace
        )
        
        if existing_entries.matches:
            return {
                "status": "duplicate",
                "message": "File already exists in vector database",
                "existing_id": existing_entries.matches[0].id
            }

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(extracted_text)
        
        metadatas = [{
            "curriculum": curriculum,
            "student_class": student_class,
            "subject": subject,
            "filename": file.filename,
            "type": "document",
            "chunk_idx": i,
            "text": chunk[:500]
        } for i, chunk in enumerate(chunks)]

        background_tasks.add_task(
            store_in_vector_db,
            chunks=chunks,
            metadatas=metadatas,
            namespace=namespace
        )

        return {
            "status": "success",
            "filename": file.filename,
            "message": "File processed and queued for vector storage",
            "processing_times": {
                "file_save": f"{time.time()-start_time:.2f}s",
                "text_extraction": f"{time.time()-extract_start:.2f}s"
            },
            "timeout_occurred": timeout_occurred,
            "chunk_count": len(chunks),
            "namespace": namespace
        }

    except HTTPException as e:
        raise e
    except Exception as e:
        logging.error(f"Upload failed: {str(e)}", exc_info=True)
        raise HTTPException(
            status_code=500,
            detail="Internal server error during file processing"
        )
    finally:
        if file_path and os.path.exists(file_path):
            try:
                os.remove(file_path)
                logging.info("Temporary file cleaned up")
            except Exception as e:
                logging.warning(f"Failed to delete temp file: {str(e)}")

def store_in_vector_db(chunks: List[str], metadatas: List[Dict], namespace: str):
    try:
        logging.info(f"Starting vector storage for {len(chunks)} chunks in namespace {namespace}")
        
        vectorstore = PineconeVectorStore.from_existing_index(
            index_name=PINECONE_INDEX_NAME,
            embedding=embedding_model,
            text_key="text",
            namespace=namespace
        )
        
        batch_size = 50 
        for i in range(0, len(chunks), batch_size):
            batch_chunks = chunks[i:i + batch_size]
            batch_metadatas = metadatas[i:i + batch_size]
            
            try:
                embeddings = embedding_model.embed_documents(batch_chunks)
                
                records = []
                for j, (chunk, metadata) in enumerate(zip(batch_chunks, batch_metadatas)):
                    record = {
                        "id": f"{namespace}_{i+j}_{uuid.uuid4().hex[:8]}",
                        "values": embeddings[j],
                        "metadata": metadata
                    }
                    records.append(record)
                
                vectorstore._index.upsert(vectors=records, namespace=namespace)
                logging.info(f"Upserted batch {i//batch_size + 1} with {len(records)} vectors")
                
            except Exception as e:
                logging.error(f"Failed to process batch {i//batch_size + 1}: {str(e)}")
                continue
        
        logging.info(f"Completed vector storage for namespace {namespace}")
        
    except Exception as e:
        logging.error(f"Vector storage failed: {str(e)}", exc_info=True)


BASE_DIR = Path("/home/ubuntu/ai_speech/ai_speech_model/ai_speech_fastapi")

class ChatResponse(TypedDict):
    question: str
    answer: Annotated[str, "Do not include URLs, only the valuable answer."]
    diagram_urls: Annotated[str, "Provide exactly one relevant diagram URL."]
    source_documents: str

model_name = OllamaLLM(model="mistral")

def parse_llm_response(response_text: str):
    """Parse the LLM response which should be in JSON format"""
    try:
        start_idx = response_text.find('{')
        end_idx = response_text.rfind('}') + 1
        json_str = response_text[start_idx:end_idx]
        return json.loads(json_str)
    except (json.JSONDecodeError, ValueError) as e:
        logging.error(f"Failed to parse LLM response: {e}")
        return {
            "question": "",
            "answer": response_text,
            "url": ""
        }   

HISTORY_DIR = "/chat_histroy_tmp"
MAX_HISTORY = 10

HISTORY_DIR = os.path.join(os.getcwd(), "chat_histroy_tmp")
os.makedirs(HISTORY_DIR, exist_ok=True)

def load_chat_history(file_path):
    """Load existing chat history or return empty list."""
    if os.path.exists(file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except json.JSONDecodeError:
            logging.warning(f"Corrupted history file: {file_path}, resetting.")
            return []
    return []

def save_chat_history(file_path, history):
    """Save chat history safely (atomic write)."""
    tmp_path = f"{file_path}.tmp"
    with open(tmp_path, "w", encoding="utf-8") as f:
        json.dump(history, f, ensure_ascii=False, indent=2)
    os.replace(tmp_path, file_path)

def cleanup_old_history(username, student_class, subject, curriculum, today_date):
    """Delete all history files for the user that are not from today."""
    pattern = f"{username}_{student_class}_{subject}_{curriculum}_*.json"
    old_files = glob.glob(os.path.join(HISTORY_DIR, pattern))
    for file_path in old_files:
        if today_date not in os.path.basename(file_path):
            try:
                os.remove(file_path)
                logging.info(f"Deleted old history file: {file_path}")
            except Exception as e:
                logging.error(f"Failed to delete old history file {file_path}: {e}")

@app.post("/chat/", response_model=ChatResponse)
async def chat(request: ChatRequest):
    try:
        today_date = datetime.now().strftime("%Y-%m-%d")
        chat_student_class = request.student_class.strip().replace(" ", "_")
        cleanup_old_history(
            request.username.strip(),
            chat_student_class,
            request.subject.strip(),
            request.curriculum.strip(),
            today_date
        )

        history_filename = (
            f"{request.username.strip()}_"
            f"{chat_student_class}_"
            f"{request.subject.strip()}_"
            f"{request.curriculum.strip()}_"
            f"{today_date}.json"
        )
        history_file_path = os.path.join(HISTORY_DIR, history_filename)

        chat_history = load_chat_history(history_file_path)

        namespace = f"{request.curriculum.strip()}_{request.student_class.strip()}_{request.subject.strip()}"
        base_dir = f"/home/ubuntu/ai_speech/ai_speech_model/ai_speech_fastapi/diagrams/{request.username.strip()}_{today_date}/extracted_images"
        image_extensions = ['*.png', '*.jpg', '*.jpeg', '*.gif', '*.bmp', '*.tiff']
        image_files = []
        for ext in image_extensions:
            image_files.extend(glob.glob(os.path.join(base_dir, ext)))

        vectorstore = PineconeVectorStore.from_existing_index(
            index_name=PINECONE_INDEX_NAME,
            embedding=embedding_model,
            text_key="text",
            namespace=namespace
        )
        retriever = vectorstore.as_retriever(search_kwargs={
            "k": 8,
            "filter": {
                "subject": request.subject.strip(),
                "curriculum": request.curriculum.strip(),
                "student_class": request.student_class.strip(),
            }
        })

        retrieved_docs = retriever.invoke(request.question.strip())
        retrieved_docs_content = " ".join(doc.page_content for doc in retrieved_docs)

        history_context = "\n".join(
            f"Q: {h['question']}\nA: {h['answer']}" for h in chat_history[-MAX_HISTORY:]
        )

        prompt_template = PromptTemplate(
            template="""
                You are an expert educator providing clear, concise answers to students.
                Use ONLY the provided context and chat history to answer the question.

                Chat History:
                {history}

                Context:
                {context}

                Available Images (use exactly one relevant images if it is matching more than 65 persentage):
                {image_files}

                Rules:
                - Always look at history first and give the answer according to that also and look at chat history as well before giving the answer any.
                - Respond in strict JSON format only
                - The answer should not contain any URLs
                - Image paths must be EXACTLY one of the provided paths
                - If no image is relevant, use empty string for url but first try to find relevant image if not found then give empty string
                - Follow this exact format:
                {{
                    "question": "the asked question",
                    "answer": "your answer without any URLs",
                    "url": "relevant image URL from above or empty string"
                }}

                Question: {question}

                Response:
            """,
            input_variables=["history","question", "image_files", "context"]
        )

        qa_chain = LLMChain(llm=model_name, prompt=prompt_template)

        result = qa_chain.invoke({
            "history": history_context,
            "question": request.question.strip(),
            "image_files": "\n".join(image_files) if image_files else "No images available",
            "context": retrieved_docs_content
        })

        parsed_response = parse_llm_response(result["text"])
        response_url = parsed_response.get("url", "")

        if response_url == "empty string" or response_url == "empty_string" or response_url == "":
            response_url = ""

        elif response_url.lstrip().startswith("/diagrams"):
            response_url = response_url

        elif response_url.lstrip().startswith("/home"):
            response_url = response_url.split("/diagrams")[1]
            response_url = f"/diagrams{response_url}"
            
        else:
            response_url = f"/diagrams/{request.username.strip()}_{today_date}/extracted_images/{response_url}.png"

        if response_url.endswith(".png.png"):
             response_url = response_url.removesuffix(".png.png") + ".png"

        full_answer = parsed_response.get("answer","")
        
        full_answer = re.sub(r"(/home\S+\.(?:png|jpg|jpeg|gif|bmp|tiff))", "", full_answer, flags=re.IGNORECASE)
        full_answer = re.sub(r"(/diagrams\S+\.(?:png|jpg|jpeg|gif|bmp|tiff))", "", full_answer, flags=re.IGNORECASE)
        full_answer = re.sub(r"\s+", " ", full_answer).strip()

        chat_history.append({"question": request.question.strip(), "answer": full_answer})
        if len(chat_history) > MAX_HISTORY:
            chat_history = chat_history[-MAX_HISTORY:]
        save_chat_history(history_file_path, chat_history)

        return ChatResponse(
            question=parsed_response.get("question", request.question.strip()),
            answer=full_answer,
            diagram_urls=response_url,
            source_documents=retrieved_docs_content
        )

    except Exception as e:
        logging.error(f"Chat error: {str(e)}", exc_info=True)
        return ChatResponse(
            question=request.question,
            answer="An error occurred while processing your request.",
            diagram_urls="",
            source_documents="",
        )


# MAX_HISTORY = 10

# def load_chat_history(file_path):
#     if os.path.exists(file_path):
#         with open(file_path, "r", encoding="utf-8") as f:
#             return json.load(f)
#     return []

# def save_chat_history(file_path, history):
#     with open(file_path, "w", encoding="utf-8") as f:
#         json.dump(history, f, ensure_ascii=False, indent=2)

        
# @app.post("/chat/", response_model=ChatResponse)
# async def chat(request: ChatRequest):
#     try:

#         history_filename = f"{request.username.strip()}_{request.student_class.strip()}_{request.subject.strip()}_{request.curriculum.strip()}.json"
#         history_file_path = os.path.join("/chat_histroy_tmp", history_filename)

#         chat_history = load_chat_history(history_file_path)

#         namespace = f"{request.curriculum.strip()}_{request.student_class.strip()}_{request.subject.strip()}"
#         today_date = datetime.now().strftime("%Y-%m-%d")
#         base_dir = f"/home/ubuntu/ai_speech/ai_speech_model/ai_speech_fastapi/diagrams/{request.username.strip()}_{today_date}/extracted_images"
        
#         image_extensions = ['*.png', '*.jpg', '*.jpeg', '*.gif', '*.bmp', '*.tiff']
#         image_files = []
#         for ext in image_extensions:
#             image_files.extend(glob.glob(os.path.join(base_dir, ext)))

#         vectorstore = PineconeVectorStore.from_existing_index(
#             index_name=PINECONE_INDEX_NAME,
#             embedding=embedding_model,
#             text_key="text",
#             namespace=namespace
#         )

#         retriever = vectorstore.as_retriever(search_kwargs={
#             "k": 8,
#             "filter": {
#                 "subject": request.subject.strip(),
#                 "curriculum": request.curriculum.strip(),
#                 "student_class": request.student_class.strip(),
#             }
#         })

#         retrieved_docs = retriever.invoke(request.question.strip())
#         retrieved_docs_content = " ".join(doc.page_content for doc in retrieved_docs)

#         history_context = "\n".join(
#             f"Q: {h['question']}\nA: {h['answer']}" for h in chat_history[-MAX_HISTORY:]
#         )

#         prompt_template = PromptTemplate(
#             template="""
#                 You are an expert educator providing clear, concise answers to students.
#                 Use ONLY the provided context and chat history to answer the question.

#                 Chat History:
#                 {history}

#                 Context:
#                 {context}

#                 Available Images (use exactly one relevant images if it is matching more than 65 persentage):
#                 {image_files}

#                 Rules:
#                 - Respond in strict JSON format only
#                 - The answer should not contain any URLs
#                 - Image paths must be EXACTLY one of the provided paths
#                 - If no image is relevant, use empty string for url but first try to find relevant image if not found then give empty string
#                 - Follow this exact format:
#                 {{
#                     "question": "the asked question",
#                     "answer": "your answer without any URLs",
#                     "url": "relevant image URL from above or empty string"
#                 }}

#                 Question: {question}

#                 Response:
#             """,
#             input_variables=["history","question", "image_files", "context"]
#         )

#         qa_chain = LLMChain(
#             llm=model_name,
#             prompt=prompt_template
#         )

#         result = qa_chain.invoke({
#             "history": history_context,
#             "question": request.question.strip(),
#             "image_files": "\n".join(image_files) if image_files else "No images available",
#             "context": retrieved_docs_content
#         })

#         parsed_response = parse_llm_response(result["text"])
#         response_url = parsed_response.get("url", "")

#         if response_url == "empty string" or response_url == "empty_string" or response_url == "":
#             response_url = ""

#         elif response_url.lstrip().startswith("/diagrams"):
#             response_url = response_url

#         elif response_url.lstrip().startswith("/home"):
#             response_url = response_url.split("/diagrams")[1]
#             response_url = f"/diagrams{response_url}"
            
 
#         else:
#             response_url = f"/diagrams/{request.username.strip()}_{today_date}/extracted_images/{response_url}.png"

#         if response_url.endswith(".png.png"):
#              response_url = response_url.removesuffix(".png.png") + ".png"

#         full_answer = parsed_response.get("answer","")
        
#         full_answer = re.sub(r"(/home\S+\.(?:png|jpg|jpeg|gif|bmp|tiff))", "", full_answer, flags=re.IGNORECASE)
#         full_answer = re.sub(r"(/diagrams\S+\.(?:png|jpg|jpeg|gif|bmp|tiff))", "", full_answer, flags=re.IGNORECASE)
#         full_answer = re.sub(r"\s+", " ", full_answer).strip()

#         chat_history.append({"question": request.question.strip(), "answer": full_answer})
#         if len(chat_history) > MAX_HISTORY:
#             chat_history = chat_history[-MAX_HISTORY:]
#         save_chat_history(history_file_path, chat_history)

#         if os.path.exists(history_file_path):
#             os.remove(history_file_path)

#         return ChatResponse(
#             question=parsed_response.get("question", request.question.strip()),
#             answer=full_answer,
#             diagram_urls=response_url,
#             source_documents=retrieved_docs_content
#         )

#     except Exception as e:
#         logging.error(f"Chat error: {str(e)}", exc_info=True)
#         return ChatResponse(
#             question=request.question,
#             answer="An error occurred while processing your request.",
#             diagram_urls="",
#             source_documents="",
#         )


@app.get("/view-image/{full_path:path}")
def view_image(full_path: str):
    image_path = (BASE_DIR / full_path.lstrip("/")).resolve()

    if not str(image_path).startswith(str(BASE_DIR)):
        raise HTTPException(status_code=400, detail="Invalid file path")

    if not image_path.exists():
        raise HTTPException(status_code=404, detail="Image not found")

    return FileResponse(image_path)



# @app.post("/chat/")
# async def chat(request: ChatRequest):
#     try:
#         question = request.question.strip()
#         subject = request.subject.strip()
#         curriculum = request.curriculum.strip()
#         student_class = request.student_class.strip()
#         username = request.username.strip()
#         namespace = f"{curriculum}_{student_class}_{subject}"
#         today_date = datetime.now().strftime("%Y-%m-%d")
#         base_dir = f"diagrams/{username}_{today_date}/extracted_images"
        

#         vectorstore = PineconeVectorStore.from_existing_index(
#             index_name=PINECONE_INDEX_NAME,
#             embedding=embedding_model,
#             text_key="text",
#             namespace=namespace
#         )

#         retriever = vectorstore.as_retriever(search_kwargs={
#             "k": 5,
#             "filter": {
#                 "subject": subject,
#                 "curriculum": curriculum,
#                 "student_class":student_class,
#             }
#         })

#         retrieved_docs = retriever.invoke(question)

#         for i, doc in enumerate(retrieved_docs, start=1):
#             print(f"\nDocument {i}:\n{doc.page_content}")

#         prompt_template = PromptTemplate.from_template("""
#         You are an expert educator providing clear, concise answers to students.
#         Extract the most relevant information to answer the question using ONLY the provided context.

#         Follow these rules:
#         1. Answer in complete, well-structured sentences.
#         2. Do not mention page numbers or document structure.
#         3. If context doesn't contain any content, say "This information is not in our materials."
#         4. Be factual and avoid speculation.
#         5. Use proper grammar and spelling.
#         6. Keep your answer concise and to the point.
#         7. Do not include '\\n' or '*' in your output.
#         8. Do not include escape characters like \\n, \\, \", \n or any slashes.
#         9. Do not use markdown symbols like '*', '-', '`', or backslashes.

#         Context: {context}
#         Question: {question}
#         Answer:
#         """)

#         qa_chain = RetrievalQA.from_chain_type(
#             llm=model_name,
#             chain_type="stuff",
#             retriever=retriever,
#             chain_type_kwargs={"prompt": prompt_template},
#             return_source_documents=True
#         )

#         result = qa_chain.invoke({"query": question})

#         return {
#             "question": question,
#             "answer": result["result"],
#             "source_documents": [doc.metadata for doc in result["source_documents"]]
#         }

#     except Exception as e:
#         logging.error(f"Chat error: {str(e)}", exc_info=True)
#         return {
#             "question": request.question,
#             "answer": "An error occurred while processing your request.",
#             "error": str(e)
#         }


@app.get("/health")
def welcome_page():
    return {"Message": "Welcome the ai speech module page."}

chat_history = []

model = ChatOllama(
            model="mistral", 
            model_kwargs={"temperature": 0.8}
        )

async def system_message(topic, mood, student_class, level) -> SystemMessage:
    parser = StrOutputParser()

    prompt_template = PromptTemplate(template="""
    You are a friendly, knowledgeable teaching assistant. Your purpose is to:
    1. Introduce the topic in a friendly manner
    2. Answer questions conversationally
    3. Never reveal internal project details
    4. Keep responses under 100 words
                                     
    Topic: {topic}
    Mood: {mood}
    Student Class: {student_class}
    Level: {level}
    5. Introduce on understanding the topic always.""",      
    input_variables=["topic", "mood", "student_class", "level"])

    chain = prompt_template | model | parser
    result = await chain.ainvoke({
        "topic": topic,
        "mood": mood,
        "student_class": student_class,
        "level": level
    })
    return result


async def initialize_essay_document(username, student_topic, student_class, mood, accent, chat_history):
    """Initialize a new essay document in Firestore"""
    try:
        serializable_chat_history = []
        for message in chat_history:
            if isinstance(message, (AIMessage, SystemMessage, HumanMessage)):
                message_dict = {
                    'type': 'chat_message',
                    'content': message.content,
                    'message_type': 'system' if isinstance(message, SystemMessage) else 
                                   'ai' if isinstance(message, AIMessage) else 
                                   'human'
                }
                serializable_chat_history.append(message_dict)
            else:
                serializable_chat_history.append(str(message))

        essay_data = {
            "username": username,
            "topic": student_topic,
            "student_class": student_class,
            "mood": mood,
            "accent": accent,
            "chat_history": serializable_chat_history,
            "created_at": datetime.now(),
            "status": "in_progress",
            "chunks": [],
            "average_scores": {}
        }

        _, doc_ref = db.collection("essays").add(essay_data)
        return doc_ref.id
        
    except Exception as e:
        logging.error(f"Failed to initialize essay document: {str(e)}", exc_info=True)
        raise


from collections import deque

feedback_queue = deque()
feedback_event = asyncio.Event()


async def feedback_processor():
    while True:
        await feedback_event.wait()
        while feedback_queue:
            websocket, feedback = feedback_queue.popleft()
            try:
                await websocket.send_json({
                    "type": "feedback",
                    "data": feedback
                })
            except Exception as e:
                logging.error(f"Failed to send feedback: {str(e)}")
        feedback_event.clear()

@app.on_event("startup")
async def startup_event():
    asyncio.create_task(feedback_processor())


async def finalize_listening_session(result, transcribed_text, websocket):
    try:
        fluency = result.get("fluency")
        pronunciation = result.get("pronunciation")
        emotion = result.get("emotion")
        grammar_score = result.get("grammar")

        prompt_template = PromptTemplate(template = """
            You are an expert English teacher providing comprehensive feedback on a student's listening comprehension and speaking response.

            STUDENT'S COMPLETE RESPONSE:
            "{transcribed_text}"

            PERFORMANCE METRICS:
            - Average Fluency: {fluency}/10
            - Average Pronunciation: {pronunciation}/10  
            - Dominant Emotion: {emotion}
            - Grammar Score: {grammar_score}/10

            DETAILED ANALYSIS REQUIRED:
            Provide comprehensive feedback in JSON format with these specific sections:

            {{
                "listening_comprehension": {{
                    "content_understanding": "How well did the student understand the main ideas? (score 1-10 with explanation)",
                    "detail_retention": "How well did they retain and mention specific details? (score 1-10 with explanation)", 
                    "key_points_covered": ["List the main points from original content that student mentioned"],
                    "missed_points": ["Important points from original content that student missed"],
                    "comprehension_score": 8.5
                }},
                "speaking_performance": {{
                    "fluency_assessment": "Detailed fluency analysis with specific examples from speech",
                    "pronunciation_assessment": "Pronunciation strengths and areas for improvement with examples",
                    "grammar_assessment": "Grammar accuracy with specific examples of errors and corrections",
                    "vocabulary_usage": "Assessment of vocabulary level and word choice appropriateness",
                    "speaking_clarity": "Overall clarity and coherence of expression"
                }},
                "technical_metrics": {{
                    "speaking_rate_analysis": "Assessment of speaking speed and pace variation",
                    "pause_analysis": "Analysis of pause placement and duration appropriateness",
                    "filler_word_usage": "Analysis of filler words and speech hesitations",
                    "prosody_evaluation": "Intonation, stress, and rhythm assessment"
                }},
                "detailed_suggestions": [
                    "Specific suggestion 1 with example",
                    "Specific suggestion 2 with example", 
                    "Specific suggestion 3 with example",
                    "Specific suggestion 4 with example",
                    "Specific suggestion 5 with example"
                ],
                "strengths": [
                    "Specific strength 1 with example",
                    "Specific strength 2 with example",
                    "Specific strength 3 with example"
                ],
                "practice_recommendations": [
                    "Specific practice activity 1",
                    "Specific practice activity 2",
                    "Specific practice activity 3"
                ],
                "overall_scores": {{
                    "fluency": {fluency},
                    "pronunciation": {pronunciation},
                    "grammar": {grammar_score},
                    "emotion": {emotion}
                }},
                "improvement_priority": "The single most important area to focus on next",
                "encouragement": "Personalized encouraging message highlighting progress and potential"
            }}

            Provide detailed, specific feedback that helps the student understand exactly what they did well and what needs improvement.
            """,
            input_variables=["transcribed_text", "fluency", "pronunciation", "emotion", "grammar_score","fluency","pronounciation","grammar_score","emotion"]
            )
        parser = StrOutputParser()
        chain = prompt_template | model | parser
        feedback = await chain.ainvoke({
            "transcribed_text": transcribed_text,
            "fluency": fluency,
            "pronunciation": pronunciation,
            "emotion": emotion,
            "grammar_score": grammar_score,
            "fluency": fluency,
            "pronunciation": pronunciation,
            "emotion": emotion,
            "grammar_score": grammar_score,
        })

        feedback_queue.append((websocket, feedback))
        feedback_event.set()
        
        return feedback

    except Exception as e:
        logging.info(f"Session finalization error: {e}")


import asyncio

async def _safe_finalize_feedback(result, text, websocket, session_state):
    try:
        feedback_result = await finalize_listening_session(result, text, websocket)
        session_state["last_feedback"] = feedback_result
        logging.info(f"this is feedback of chunk : {feedback_result}")
    except Exception:
        logging.exception("finalize_listening_session failed")




TEMP_DIR = os.path.abspath("temp_chunks")
os.makedirs(TEMP_DIR, exist_ok=True) 

@app.websocket("/ws/assistant")
async def audio_ws_assistant(websocket: WebSocket):
    await websocket.accept()
    query_params = parse_qs(websocket.url.query)
    username = query_params.get("username", [None])[0]
    token = query_params.get("token", [None])[0]
    student_topic = query_params.get("topic", [None])[0]
    student_class = query_params.get("student_class", [None])[0]
    mood = query_params.get("mood", [None])[0]
    accent = query_params.get("accent", [None])[0]

    # Initialize audio processor
    audio_processor = AdvancedAudioProcessor()
    
    chat_history = []
    essay_id = await initialize_essay_document(
        username=username,
        student_topic=student_topic,
        student_class=student_class,
        mood=mood,
        accent=accent,
        chat_history=chat_history
    )

    logging.info(f"essay_id : {essay_id}")
    
    try:
        await websocket.send_json({"action": "essay_initialized", "essay_id": essay_id})
    except Exception as e:
        logging.error(f"Failed to send initial essay_id: {str(e)}")

    ai_response = await system_message(student_topic, mood, student_class, accent)
    chat_history.append(SystemMessage(content=ai_response))
    
    if not username or not token:
        await websocket.close(code=4001)
        logging.info("Username or token missing.")
        return

    logging.info(f"[WS] Authenticated connection from {username}")
    
    config = {
        'silence_threshold': 2,
        'min_utterance_length': 4,
        'min_speech_duration': 2.0,
        'silence_dBFS_threshold': -50,
        'processing_cooldown': 2.5,
        'blocklist': [
            'you', 'thank you', 'tchau', 'thanks', 'ok', 'Obrigado.', 'E aí', '',
            'me', 'hello', 'hi', 'hey', 'okay', 'thanks', 'thank', 'obrigado',
            'tchau.', 'bye', 'goodbye', 'me.', 'you.', 'thank you.',"I'm going to take a picture of the sea","Kansai International Airport",
            "Thank you for watching!","1 tbsp of salt",'Teksting av Nicolai Winther',
            'ん'
        ],
        'max_repetitions': 2,
        'max_silence': 15.0,
        'chunk_duration': 0.5,
    }

    session_state = {
        'silvero_model': True,
        'audio_buffer': AudioSegment.empty(),
        'speech_buffer': AudioSegment.empty(),
        'text_buffer': [],
        'last_speech_time': time.time(),
        'conversation_active': False,
        'processing_active': False,
        'assistant_speaking': False,
        'chunk_index': 0,
        'chunk_results': [],
        'text_output': [],
        'has_speech': False,
        'last_processing_time': 0,
        'consecutive_silence_chunks': 0,
        'active_speech_duration': 0.0,
        'current_noise_level': None
    }
    
    topic = Topic()
    session_temp_dir = tempfile.mkdtemp(prefix=f"{username}_", dir=TEMP_DIR)

    session_state["assistant_speaking"] = True
    response_audio = await topic.text_to_speech_assistant(ai_response, username, session_temp_dir)
    sleep_time = await send_audio_response(websocket, response_audio)
    await asyncio.sleep(sleep_time + 2)
    session_state["assistant_speaking"] = False

    logging.info(f"[Session] Temp dir created at {session_temp_dir}")
    
    final_output = os.path.join(session_temp_dir, f"{username}_output.wav")
    transcript_path = os.path.join(session_temp_dir, f"{username}_transcript.txt")

    try:
        while True:
            message = await websocket.receive()

            if message["type"] == "websocket.disconnect":
                logging.info(f"[WS] {username} disconnected.")
                break

            if message["type"] == "websocket.receive" and "bytes" in message:
                if (session_state['processing_active'] or 
                    session_state['assistant_speaking'] or
                    (time.time() - session_state['last_processing_time'] < config['processing_cooldown'])):
                    continue    

                current_time = time.time()
                
                # Process audio with advanced noise cancellation
                processed_bytes = await audio_processor.process_audio_chunk(message["bytes"])
                
                new_chunk = AudioSegment(
                    data=processed_bytes,
                    sample_width=2,
                    frame_rate=16000,
                    channels=1
                )


                logging.info(f"new chunks data ------------>>>{new_chunk}")
                
                # Update session state with current noise level
                session_state['current_noise_level'] = audio_processor.last_noise_level
                
                
                loudness = new_chunk.dBFS
                is_silent = loudness < config['silence_dBFS_threshold']
                session_state['audio_buffer'] += new_chunk
                
                if is_silent:
                    session_state['consecutive_silence_chunks'] += 1
                    logging.info(f"🟡 Silent chunk detected (loudness: {loudness:.2f} dB)")
                    
                    if (session_state['has_speech'] and 
                        session_state['consecutive_silence_chunks'] * config['chunk_duration'] >= config['silence_threshold'] and
                        session_state['active_speech_duration'] >= config['min_speech_duration']):
                        if not session_state['assistant_speaking']:
                            await process_buffered_audio(
                                final_output, transcript_path, session_state, 
                                websocket, username, session_temp_dir, topic, 
                                config, student_topic, student_class, mood, accent, chat_history
                            )
                else:
                    if current_time - session_state['last_processing_time'] < config['processing_cooldown']:
                        continue
                        
                    session_state['consecutive_silence_chunks'] = 0
                    session_state['has_speech'] = True
                    session_state['speech_buffer'] += new_chunk
                    session_state['active_speech_duration'] += config['chunk_duration']
                    session_state['last_speech_time'] = current_time
                    logging.info(f"🔵 Speech chunk detected (loudness: {loudness:.2f} dB), active duration: {session_state['active_speech_duration']:.2f}s")
                    
                    chunk_filename = os.path.join(session_temp_dir, f"chunk_temp_{session_state['chunk_index']}.wav")
                    new_chunk.export(chunk_filename, format="wav")
                    session_state['chunk_index'] += 1
                    
                    vad_result = await topic.silvero_vad(chunk_filename)
                    current_silence = vad_result.get("duration", 0.0)
                    
                    if current_silence > 0.3:
                        logging.info(f"VAD detected silence gap: {current_silence:.2f}s")
                        
                        if (current_silence >= config['silence_threshold']):
                            
                            if not session_state['assistant_speaking']:
                                await process_buffered_audio(
                                    final_output, transcript_path, session_state, 
                                    websocket, username, session_temp_dir, topic, 
                                    config, student_topic, student_class, mood, accent, chat_history
                                )

    except WebSocketDisconnect:
        logging.warning(f"[WS] {username} disconnected unexpectedly")
    except Exception as e:
        logging.error(f"Unexpected error: {str(e)}", exc_info=True)
    finally:
        essay_id = finalize_session(session_state, username, session_temp_dir, topic, essay_id, chat_history)

async def process_buffered_audio(
    final_output: str,
    transcript_path: str,
    session_state: dict,
    websocket: WebSocket,
    username: str,
    temp_dir: str,
    topic: Topic,
    config: dict,
    student_topic: str,
    student_class: str,
    mood: str,
    accent: str,
    chat_history: list
):
    if len(session_state['audio_buffer']) == 0:
        return
        
    session_state['processing_active'] = True
    try:
        buffer_filename = os.path.join(temp_dir, f"buffered_{time.time()}.wav")
        session_state['audio_buffer'].export(buffer_filename, format="wav")
        
        transcribed_text = await topic.speech_to_text(buffer_filename,username)
        grammar_score = await topic.grammar_checking(transcribed_text)
        logging.info(f"Transcribed: {transcribed_text}")
        session_state['text_buffer'].append(transcribed_text)
        session_state['text_output'].append(transcribed_text)
        
        if os.path.exists(final_output):
            existing_audio = AudioSegment.from_wav(final_output)
            combined_audio = existing_audio + session_state['audio_buffer']
        else:
            combined_audio = session_state['audio_buffer']
        combined_audio.export(final_output, format="wav")
        
        with open(transcript_path, "w", encoding="utf-8") as f:
            f.write(" ".join(session_state['text_output']).strip())
        
        clean_text = transcribed_text.lower().strip()
        if (len(transcribed_text.split()) >= config['min_utterance_length'] and 
            clean_text not in config['blocklist']):
            
            async with aiohttp.ClientSession() as session:
                with open(buffer_filename, "rb") as f:
                    form = aiohttp.FormData()
                    form.add_field("file", f, filename=os.path.basename(buffer_filename), content_type="audio/wav")
                    async with session.post(f"{CPU_API_BASE}/detect-emotion", data=form) as res:
                        emotion_data = await res.json()
                        emotion = emotion_data.get("emotion")

                async with session.post(
                    f"{CPU_API_BASE}/fluency-score",
                    json={"text": transcribed_text}
                ) as res:
                    fluency_data = await res.json()
                    fluency = fluency_data.get("fluency")

                with open(buffer_filename, "rb") as f:
                    form = aiohttp.FormData()
                    form.add_field("file", f, filename=os.path.basename(buffer_filename), content_type="audio/wav")
                    async with session.post(f"{CPU_API_BASE}/pronunciation-score", data=form) as res:
                        pron_data = await res.json()
                        pronunciation = pron_data.get("pronunciation")

                result = {
                    "fluency": fluency,
                    "pronunciation": pronunciation,
                    "emotion": emotion,
                    "grammar": grammar_score
                }

                # asyncio.create_task(
                #     _safe_finalize_feedback(result, transcribed_text, websocket, session_state)
                # )

                if None in result.values():
                    logging.warning(f"Incomplete analysis results: {result}")

                await process_user_utterance(
                    transcribed_text, result["emotion"], result["fluency"], result["pronunciation"],
                    session_state, buffer_filename, websocket, 
                    username, temp_dir, topic, student_topic, student_class, mood, accent, chat_history
                )

        session_state['audio_buffer'] = AudioSegment.empty()
        session_state['text_buffer'] = []
        session_state['conversation_active'] = True
        
    except Exception as e:
        logging.error(f"Error processing buffered audio: {str(e)}", exc_info=True)
    finally:
        session_state['processing_active'] = False

async def process_user_utterance(
    text: str,
    emotion: str,
    fluency: float,
    pronunciation: float,
    session_state: dict,
    chunk_filename: str,
    websocket: WebSocket,
    username: str,
    session_temp_dir: str,
    topic: Topic,
    student_topic: str,
    student_class: str,
    mood: str,
    accent: str,
    chat_history: list
):
    topic.update_realtime_stats(fluency, pronunciation, emotion)

    session_state['text_output'].append(text)
    session_state['chunk_results'].append({
        "chunk_index": session_state['chunk_index'],
        "text": text,
        "emotion": emotion,
        "fluency": fluency,
        "pronunciation": pronunciation,
        "file_path": chunk_filename
    })
    session_state['chunk_index'] += 1
    scraped_data = await scraping(student_topic)

    try:
        session_state["assistant_speaking"] = True
        prompt_template = PromptTemplate(
            template="""
            ROLE: You are a friendly, knowledgeable assistant and communicate like a friend. Your purpose is to:
            1. Answer questions conversationally
            2. Never reveal internal project details
            3. Keep responses under 200 words
            5. if by chance question is much more specific topic is needed and you does not have correct answer then used the function await scraping(your_specific_topic), which help in give updated answer.
            6. If you did not have updated data then always used the {scraped_data}, for update yourself, if by chance is question is for related to udpated data then also you able to send the request on the funtion await scraping(topic_you_want_to_select)                           
            Chat History:
            {chat_history}
            considering history as well before answering the question.

            Topic: {student_topic}
            Mood: {mood}
            Student Class: {student_class}
            Level: {level}
                                         
            - for reply anything consider always the Topic, Mood , student_class, Level and history.
            - And should be reply in the way and look like a friendly, knowledgeable teaching assistant.
            - Do not mention any technical details, model architecture, or team members.
            - Focus on providing helpful, concise answers.
            - If the question is not related to the topic, politely tell that quesiton is not related to the toic                     
                                                                

            RULES:
            - NEVER mention:
            * Model architecture/type
            * Team members/credentials
            * Code implementation
            * Technical specifications
            - Always redirect technical questions to general knowledge                                                                                                 

            USER QUESTION: {question}

            RESPONSE FORMAT:
            [Answer Concise 1-10 sentence response] 
            [if needed  Optional follow-up question to continue conversation]

            EXAMPLE:
            This exaple is for you never asked that. 
            User: What model are you using?
            I focus on helping with learning concepts rather than technical details. 
            Would you like me to explain how these systems generally work?

            Current response should be:
            """,
            input_variables=["scraped_data", "question", "student_topic", "student_class", "mood", "accent", "chat_history"]
        )
        
        model = ChatOllama(
            model="mistral",
            model_kwargs={"temperature": 0.8}
        )

        parser = StrOutputParser()
        chain = prompt_template | model | parser

        ai_response = await chain.ainvoke({
            "scraped_data": scraped_data,
            "student_topic": student_topic,
            "mood": mood,
            "student_class": student_class,
            "level": accent,
            "question": text,
            "chat_history": chat_history
        })
    
        chat_history.append(AIMessage(content=ai_response))
        logging.info(f"[AI Response]: {ai_response}")

        chunk_result = {
            "chunk_index": session_state['chunk_index'],
            "text": text,
            "emotion": emotion,
            "fluency": fluency,
            "pronunciation": pronunciation,
            "file_path": chunk_filename,
            "chat_history": chat_history,
        }
        
        session_state["chunk_results"].append(chunk_result)

        response_audio = await topic.text_to_speech_assistant(ai_response, username, session_temp_dir)
        sleep_time = await send_audio_response(websocket, response_audio)
        await asyncio.sleep(sleep_time) 


    except Exception as e:
        logging.error(f"QA Error: {str(e)}")
        await send_default_response(websocket, username, session_temp_dir, topic)

    finally:
        session_state["assistant_speaking"] = False


async def scraping(topic: str) -> str:
    url = f"https://en.wikipedia.org/wiki/{topic}"
    api_endpoint = f"https://api.scrapingdog.com/scrape?api_key={scraping_api_key}&url={url}"

    text = ""
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(api_endpoint) as response:
                if response.status == 200:
                    html = await response.text()
                    soup = BeautifulSoup(html, "html.parser")
                    for tag in soup(["script", "style", "noscript"]):
                        tag.decompose()
                    text = soup.get_text(separator="\n", strip=True)
                else:
                    logging.warning(f"[SCRAPING] Failed with status {response.status}: {await response.text()}")
    except Exception as e:
        logging.exception(f"[SCRAPING ERROR] Failed to scrape data: {e}")

    return text

async def send_audio_response(websocket: WebSocket, audio_file: str) -> float:
    try:
        if not os.path.exists(audio_file):
            logging.error(f"Audio file not found: {audio_file}")
            return 0
            
        audio = AudioSegment.from_wav(audio_file)
        duration_ms = len(audio)
        
        if websocket.client_state == WebSocketState.DISCONNECTED:
            logging.warning("WebSocket disconnected before sending audio")
            return 0
            
        with open(audio_file, "rb") as f:
            await websocket.send_bytes(f.read())

        return duration_ms / 1000
        
    except Exception as e:
        logging.error(f"Failed to send audio response: {str(e)}")
        return 0

async def send_default_response(websocket: WebSocket, username: str, session_temp_dir: str, topic: Topic):
    try:
        default_response = "I didn't quite catch that. Could you please repeat?"
        response_audio = await topic.text_to_speech_assistant(default_response, username, session_temp_dir)
        await send_audio_response(websocket, response_audio)
    except Exception as e:
        logging.error(f"Failed to send default response: {str(e)}")

def cleanup_temp_files(session_temp_dir: str, chunk_results: list):
    try:
        for chunk in chunk_results:
            try:
                if os.path.exists(chunk.get('file_path', '')):
                    os.remove(chunk['file_path'])
            except Exception as e:
                logging.warning(f"Failed to delete chunk file {chunk.get('file_path', '')}: {e}")

        try:
            if os.path.exists(session_temp_dir):
                if not os.listdir(session_temp_dir):
                    os.rmdir(session_temp_dir)
        except Exception as e:
            logging.warning(f"Failed to remove temp directory {session_temp_dir}: {e}")

    except Exception as e:
        logging.error(f"Error during temp file cleanup: {str(e)}")

def finalize_session(
    session_state: dict,
    username: str,
    session_temp_dir: str,
    topic: Topic,
    essay_id: str,
    chat_history: list
) -> str:
    try:
        serializable_chat_history = []
        for message in chat_history:
            if isinstance(message, (AIMessage, SystemMessage, HumanMessage)):
                message_dict = {
                    'type': 'chat_message',
                    'content': message.content,
                    'message_type': 'system' if isinstance(message, SystemMessage) else 
                                   'ai' if isinstance(message, AIMessage) else 
                                   'human'
                }
                serializable_chat_history.append(message_dict)
            elif isinstance(message, dict):
                serializable_chat_history.append(message)
            else:
                serializable_chat_history.append(str(message))

        chunk_results_clean = []
        for chunk in session_state['chunk_results']:
            clean_chunk = {
                'chunk_index': chunk.get('chunk_index'),
                'text': chunk.get('text'),
                'emotion': chunk.get('emotion'),
                'fluency': chunk.get('fluency'),
                'pronunciation': chunk.get('pronunciation'),
            }
            chunk_results_clean.append(clean_chunk)

        essay_ref = db.collection("essays").document(essay_id)
        
        update_data = {
            "chat_history": serializable_chat_history,
            "chunks": chunk_results_clean,
            "average_scores": topic.get_average_realtime_scores(),
            "updated_at": datetime.now(),
            "status": "completed",
            "transcript": " ".join(session_state['text_output']).strip()
        }

        essay_ref.update(update_data)
        logging.info(f"Successfully updated essay document {essay_id}")
        
    except Exception as e:
        logging.error(f"Failed to finalize essay document: {str(e)}", exc_info=True)
        raise
    
    finally:
        cleanup_temp_files(session_temp_dir, session_state['chunk_results'])
    
    return essay_id


